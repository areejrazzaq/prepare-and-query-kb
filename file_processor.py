import io
import logging
from typing import List, Dict, Any
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_content: bytes) -> str:
    """
    Extract text from PDF file content.
    
    Args:
        file_content: PDF file content as bytes
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If PDF extraction fails
    """
    try:
        pdf_file = io.BytesIO(file_content)
        reader = PdfReader(pdf_file)
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}", exc_info=True)
        raise Exception(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_txt(file_content: bytes) -> str:
    """
    Extract text from TXT file content.
    
    Args:
        file_content: TXT file content as bytes
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If text extraction fails
    """
    try:
        # Try UTF-8 first, fallback to latin-1 if that fails
        try:
            text = file_content.decode('utf-8')
        except UnicodeDecodeError:
            text = file_content.decode('latin-1', errors='ignore')
        return text
    except Exception as e:
        logger.error(f"Error extracting text from TXT file: {str(e)}", exc_info=True)
        raise Exception(f"Failed to extract text from TXT file: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    """
    Extract text from DOCX file content.
    
    Args:
        file_content: DOCX file content as bytes
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If DOCX extraction fails
    """
    try:
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}", exc_info=True)
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_ppt(file_content: bytes) -> str:
    """
    Extract text from PPT/PPTX file content.
    
    Args:
        file_content: PPT/PPTX file content as bytes
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If PPT extraction fails
    """
    try:
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PPT: {str(e)}", exc_info=True)
        raise Exception(f"Failed to extract text from PPT: {str(e)}")


def extract_text(file_content: bytes, file_extension: str) -> str:
    """
    Main dispatcher function to extract text based on file extension.
    
    Args:
        file_content: File content as bytes
        file_extension: File extension (e.g., 'pdf', 'txt', 'docx', 'ppt')
        
    Returns:
        Extracted text as string
        
    Raises:
        ValueError: If file extension is not supported
        Exception: If text extraction fails
    """
    extension = file_extension.lower().lstrip('.')
    
    if extension == 'pdf':
        return extract_text_from_pdf(file_content)
    elif extension == 'txt':
        return extract_text_from_txt(file_content)
    elif extension in ['docx', 'docs']:
        return extract_text_from_docx(file_content)
    elif extension in ['ppt', 'pptx']:
        return extract_text_from_ppt(file_content)
    else:
        raise ValueError(f"Unsupported file extension: {extension}")


def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 100) -> List[str]:
    """
    Split text into fixed-size chunks with optional overlap.
    
    Args:
        text: Text to chunk
        chunk_size: Size of each chunk in characters (default: 2000)
        overlap: Number of characters to overlap between chunks (default: 100)
        
    Returns:
        List of text chunks
    """
    if not text or len(text.strip()) == 0:
        return []
    
    # Clean and normalize text
    text = text.strip()
    
    # If text is shorter than chunk_size, return as single chunk
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        # Calculate end position
        end = start + chunk_size
        
        # Extract chunk
        chunk = text[start:end]
        
        # If not the last chunk and there's more text, try to break at word boundary
        if end < len(text) and overlap > 0:
            # Look for a good break point (space, newline, or punctuation)
            break_point = chunk.rfind('\n')
            if break_point == -1 or break_point < chunk_size - 200:
                break_point = chunk.rfind('. ')
            if break_point == -1 or break_point < chunk_size - 200:
                break_point = chunk.rfind(' ')
            
            if break_point > chunk_size // 2:  # Only use break point if it's reasonable
                chunk = chunk[:break_point]
                end = start + break_point
        
        chunks.append(chunk.strip())
        
        # Move start position forward, accounting for overlap
        start = end - overlap
        if start >= len(text):
            break
    
    return chunks


def format_records_for_pinecone(chunks: List[str], assistant: str, base_id: str = None) -> List[Dict[str, Any]]:
    """
    Format chunks into Pinecone record format matching recursive_crawl.py structure.
    
    Args:
        chunks: List of text chunks
        assistant: Assistant identifier for filtering
        base_id: Base identifier for records (default: assistant name)
        
    Returns:
        List of dictionaries formatted for Pinecone upsert
    """
    if base_id is None:
        base_id = assistant
    
    records = []
    for i, chunk in enumerate(chunks, start=1):
        record = {
            '_id': f"{base_id}_{i}",
            'text': chunk,
            'assistant': assistant
        }
        records.append(record)
    
    return records

